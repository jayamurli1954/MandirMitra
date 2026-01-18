from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Configuration ---
# Colors
COLOR_PRIMARY_SAFFRON = RGBColor(255, 153, 51)  # #FF9933
COLOR_DARK_SAFFRON = RGBColor(230, 126, 34)     # #E67E22
COLOR_TEMPLE_GREEN = RGBColor(19, 136, 8)       # #138808
COLOR_DARK_GREEN = RGBColor(11, 92, 5)          # #0B5C05
COLOR_DEEP_GRAY = RGBColor(51, 51, 51)          # #333333
COLOR_LIGHT_BG = RGBColor(247, 247, 247)        # #F7F7F7
COLOR_WHITE = RGBColor(255, 255, 255)

# Fonts
FONT_HEADING = "Montserrat" # Will fallback to Arial if not installed
FONT_BODY = "Roboto"        # Will fallback to Calibri if not installed

def create_presentation():
    prs = Presentation()
    
    # 1. Setup Slide Master / Layouts (Basic simulation via code)
    # We will format elements manually for each slide to ensure strict adherence to spec
    
    # --- Helper: Create Title Slide ---
    def add_title_slide(prs, title_text, subtitle_text, footer_text):
        slide_layout = prs.slide_layouts[6] # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Background (Simulated Gradient with solid fill for python-pptx simplicity, 
        # or use an image if available. We will use Primary Saffron)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_PRIMARY_SAFFRON
        
        # Title
        title = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1.5))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_HEADING
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        sub = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(1))
        tf_sub = sub.text_frame
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle_text
        p_sub.font.name = FONT_BODY
        p_sub.font.size = Pt(24)
        p_sub.font.color.rgb = COLOR_WHITE
        p_sub.alignment = PP_ALIGN.CENTER
        
        # Video Placeholder Circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(1), Inches(2), Inches(2))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_WHITE
        circle.line.color.rgb = COLOR_WHITE
        
        vid_text = slide.shapes.add_textbox(Inches(3.5), Inches(1.8), Inches(3), Inches(1))
        tf_v = vid_text.text_frame
        p_v = tf_v.paragraphs[0]
        p_v.text = "Insert Video Here\n(MandirMitra-logo.mp4)"
        p_v.font.size = Pt(10)
        p_v.alignment = PP_ALIGN.CENTER
        
        # Footer
        ft = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        tf_ft = ft.text_frame
        p_ft = tf_ft.paragraphs[0]
        p_ft.text = footer_text
        p_ft.font.italic = True
        p_ft.font.color.rgb = COLOR_WHITE
        p_ft.alignment = PP_ALIGN.CENTER

    # --- Helper: Create Content Slide ---
    def add_content_slide(prs, title_text, bullets_list, highlight_text=None, bottom_note=None):
        slide_layout = prs.slide_layouts[6] # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Top Bar (Strip)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_PRIMARY_SAFFRON
        bar.line.fill.background()
        
        # Logo Corner
        logo = slide.shapes.add_textbox(Inches(8), Inches(0.2), Inches(2), Inches(0.5))
        tf_logo = logo.text_frame
        p_logo = tf_logo.paragraphs[0]
        p_logo.text = "MandirMitra"
        p_logo.font.name = FONT_HEADING
        p_logo.font.bold = True
        p_logo.font.size = Pt(14)
        p_logo.font.color.rgb = COLOR_DARK_SAFFRON
        p_logo.alignment = PP_ALIGN.RIGHT

        # Title
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_HEADING
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK_SAFFRON
        
        # Bullets
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
        tf_body = body.text_frame
        tf_body.word_wrap = True
        
        for bullet in bullets_list:
            p = tf_body.add_paragraph()
            p.text = bullet
            p.font.name = FONT_BODY
            p.font.size = Pt(20)
            p.font.color.rgb = COLOR_DEEP_GRAY
            p.level = 0
            p.space_after = Pt(10)
            
        # Highlight/Bottom Note
        if highlight_text:
            box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.8))
            tf_h = box.text_frame
            p_h = tf_h.paragraphs[0]
            p_h.text = highlight_text
            p_h.font.bold = True
            p_h.font.color.rgb = COLOR_TEMPLE_GREEN
            p_h.font.size = Pt(18)
            
        if bottom_note:
             # Add note as small text at bottom
             pass 

    # --- Helper: Two Column Slide ---
    def add_two_col_slide(prs, title_text, left_title, left_bullets, right_title, right_bullets):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Top Bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_PRIMARY_SAFFRON
        bar.line.fill.background()
        
        # Title
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_HEADING
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK_SAFFRON

        # Left Col
        l_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(4.25), Inches(4.5))
        l_box.fill.solid()
        l_box.fill.fore_color.rgb = RGBColor(255, 243, 224) # Light Orange
        l_box.line.fill.background()
        
        l_title = slide.shapes.add_textbox(Inches(0.7), Inches(2), Inches(3.8), Inches(0.5))
        l_title.text_frame.text = left_title
        l_title.text_frame.paragraphs[0].font.bold = True
        l_title.text_frame.paragraphs[0].font.color.rgb = COLOR_DARK_SAFFRON
        
        l_body = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(3.8), Inches(3.5))
        tf_l = l_body.text_frame
        tf_l.word_wrap = True
        for b in left_bullets:
            p = tf_l.add_paragraph()
            p.text = b
            p.font.size = Pt(16)
            p.font.color.rgb = COLOR_DEEP_GRAY
            
        # Right Col
        r_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.25), Inches(1.8), Inches(4.25), Inches(4.5))
        r_box.fill.solid()
        r_box.fill.fore_color.rgb = RGBColor(232, 245, 233) # Light Green
        r_box.line.fill.background()

        r_title = slide.shapes.add_textbox(Inches(5.45), Inches(2), Inches(3.8), Inches(0.5))
        r_title.text_frame.text = right_title
        r_title.text_frame.paragraphs[0].font.bold = True
        r_title.text_frame.paragraphs[0].font.color.rgb = COLOR_DARK_GREEN

        r_body = slide.shapes.add_textbox(Inches(5.45), Inches(2.5), Inches(3.8), Inches(3.5))
        tf_r = r_body.text_frame
        tf_r.word_wrap = True
        for b in right_bullets:
            p = tf_r.add_paragraph()
            p.text = b
            p.font.size = Pt(16)
            p.font.color.rgb = COLOR_DEEP_GRAY

    # --- SLIDE 1: Title ---
    add_title_slide(prs, 
                    "MandirMitra", 
                    "Integrated Temple Management Software\nSimplifying Seva, Accounting & Administration",
                    "\"Sacred workflows, simplified.\"")

    # --- SLIDE 2: Overview ---
    add_content_slide(prs, "What is MandirMitra?", [
        "Purpose-built temple management software for Indian temples of all sizes.",
        "Integrates Donations, Sevas, Accounting, Inventory, Assets & Panchang in one system.",
        "Designed for transparency, audit-readiness, and operational efficiency.",
        "Available as SaaS (cloud) or Standalone (on-premise) deployment."
    ])

    # --- SLIDE 3: Pain Points ---
    add_content_slide(prs, "Challenges in Temple Administration", [
        "Fragmented records: Manual ledgers, Excel sheets, loose parchis.",
        "Lack of audit trail for donations, sevas, inventory, and assets.",
        "Difficulty reconciling donation counters, UPI collections, and bank.",
        "No integrated view of seva bookings, CWIP, assets, and depreciation.",
        "Time-consuming preparation of reports for trustees, donors, auditors."
    ], "Bottom line: MandirMitra converts all of this into a single, reliable system.")

    # --- SLIDE 4: Modules (Using Bullets for simplicity in code) ---
    add_content_slide(prs, "MandirMitra – Major Modules", [
        "Devotees & Donations (Mobile-first entry)",
        "Sevas & Token Sevas (Pre-printed tokens, reconciliation)",
        "Accounting & Reports (Double-entry, Auditors view)",
        "Inventory Management (Store wise, consumption)",
        "Asset Management & Depreciation (CWIP, Gold, Land)",
        "Panchang & (Optional) Kundli"
    ])

    # --- SLIDE 5: Devotees ---
    add_content_slide(prs, "Devotee & Donation Management", [
        "Devotee master with mobile-first search and auto-fill.",
        "Fast Quick Donation Entry from dashboard.",
        "Auto-fill address from PIN code (city, state, country).",
        "Consolidated accounting: 4100 – Donation Income (Main).",
        "Multiple modes: cash, UPI, bank, with clear tracking."
    ], "Highlight: Reduces counter time and improves data quality.")

    # --- SLIDE 6: Sevas ---
    add_content_slide(prs, "Seva Management & Token Sevas", [
        "Seva master: timings, capacity, approvals, priests, pricing.",
        "Seva booking with mobile-based devotee lookup & PIN code auto-fill.",
        "Token Seva system: Color-coded pre-printed tokens for small-value sevas.",
        "Controlled serial numbers and reconciliation.",
        "Accounting consolidation: 4200 – Sevas – Main."
    ], "Emphasize internal control and audit-friendliness.")

    # --- SLIDE 7: Accounting ---
    add_content_slide(prs, "Integrated Accounting Engine", [
        "Full double-entry bookkeeping with Chart of Accounts.",
        "Automatic journal entries for Donations, Sevas, Inventory, Assets.",
        "Handling of CWIP capitalization, depreciation, revaluation.",
        "Trial Balance shows short list of main accounts with drill-down.",
        "Designed with auditors in mind – traceable, consistent, exportable."
    ])

    # --- SLIDE 8: Inventory ---
    add_content_slide(prs, "Inventory & Store Management", [
        "Item Master (Pooja materials, groceries, cleaning, etc.).",
        "Store Master (Main store, kitchen, pooja room, etc.).",
        "Purchase Entry: Dr Inventory / Cr Cash/Bank.",
        "Issue Entry: Dr Expense / Cr Inventory (consumption).",
        "Stock Report: item/store-wise quantities, values, low-stock."
    ])

    # --- SLIDE 9: Assets ---
    add_content_slide(prs, "Asset Management – Foundations", [
        "Asset Master: land, buildings, vehicles, equipment, gold/silver.",
        "Asset categorization with account mapping (1500–1999 series).",
        "Tracks original cost, book value, accumulated depreciation.",
        "Handles Fixed Assets, Precious Assets and CWIP."
    ], "Built specifically for temple asset realities (land, gold, silver).")

    # --- SLIDE 10: CWIP ---
    add_content_slide(prs, "CWIP & Procurement", [
        "CWIP projects for construction, renovation, infrastructure.",
        "Record all project expenses – Materials, Labour, Overheads.",
        "Accounting: Dr CWIP account, Cr Cash/Bank/Payables.",
        "Capitalize to Fixed Asset: Dr Asset, Cr CWIP (automatic entry).",
        "Optional tender process designed for large temples."
    ])

    # --- SLIDE 11: Depreciation ---
    add_content_slide(prs, "Flexible Depreciation – 8 Methods", [
        "Supports all major methods: Straight-Line, WDV, etc.",
        "Temple admin selects method in consultation with auditor.",
        "Automated schedules and journal entries.",
        "Dr Depreciation Expense, Cr Accumulated Depreciation.",
        "Compliant with standard accounting practices (AS/Ind AS)."
    ])

    # --- SLIDE 12: Revaluation ---
    add_content_slide(prs, "Revaluation & Disposal of Assets", [
        "Revaluation for land, buildings, gold, silver.",
        "Handles upward and downward revaluations with Reserve.",
        "Disposal types: sale, scrap, donation, loss, theft.",
        "Automatic gain/loss calculation and posting.",
        "Example: Dr Asset, Cr Revaluation Reserve."
    ], "Keeps temple balance sheet realistic and audit-ready.")

    # --- SLIDE 13: Panchang ---
    add_content_slide(prs, "Daily Panchang & Optional Kundli", [
        "Daily Panchang integrated with temple location.",
        "Configurable display settings.",
        "Kundli module: currently flagged/hidden until perfected.",
        "Only enabled on demand for users who need it."
    ])

    # --- SLIDE 14: SaaS vs Standalone (Two Col) ---
    add_two_col_slide(prs, "Deployment Options", 
                      "SaaS (Cloud)", 
                      ["Accessible from anywhere.", "Automatic backups & updates.", "Lower upfront hardware cost.", "Best for reliable internet."],
                      "Standalone (On-Premise)",
                      ["Runs on temple’s own server.", "Data stays locally.", "Temple controls backups.", "Can upgrade to SaaS later."])

    # --- SLIDE 15: Tech Stack ---
    add_content_slide(prs, "Technical Overview", [
        "Frontend: React + Material-UI with responsive design.",
        "Backend: FastAPI (Python), RESTful APIs.",
        "Database: PostgreSQL with migrations and scripts.",
        "Clean separation of Models, APIs, and Accounting engine.",
        "Robust logging and schema management."
    ])

    # --- SLIDE 16: Security ---
    add_content_slide(prs, "Security & Controls", [
        "User login & role-based access (admin, accountant, counter).",
        "Controlled token seva inventory and reconciliation.",
        "Every transaction linked to journal entries where applicable.",
        "Exportable reports for auditors and trustees.",
        "Logs and database migrations to avoid schema drift."
    ])

    # --- SLIDE 17: Benefits ---
    add_content_slide(prs, "Key Benefits to Your Temple", [
        "Single source of truth for donations, sevas, inventory, and assets.",
        "Faster counter operations with mobile-first and PIN code auto-fill.",
        "Stronger financial discipline and easier audits.",
        "Better transparency for devotees, trustees, and regulators.",
        "Scales from small local temples to large institutions."
    ], "\"Less time on administration, more time on Seva.\"")

    # --- SLIDE 18: Final ---
    add_content_slide(prs, "Final Word & Next Steps", [
        "MandirMitra combines tradition + technology.",
        "Covers: Donations > Sevas > Inventory > Assets > Accounting.",
        "Flexible SaaS or Standalone deployment.",
        "Next steps: Demo with sample data, Auditor review, Phased rollout."
    ], "Contact: contact@MandirMitra.com | Dhanyavaad")

    # Save
    prs.save('MandirMitra_Presentation.pptx')
    print("Presentation created successfully: MandirMitra_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()